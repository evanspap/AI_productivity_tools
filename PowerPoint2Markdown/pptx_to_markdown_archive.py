#!/usr/bin/env python
"""
Author: Evangelos Papadopoulos
Created: 2026-05-13
Purpose: Convert PowerPoint presentations into Markdown plus extracted images for AI/Codex/LLM-readable research archives.
License: Personal research/workflow utility.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path, PurePosixPath
from typing import Iterable
from xml.etree import ElementTree as ET


REL_NS = "{http://schemas.openxmlformats.org/package/2006/relationships}"
MEDIA_REL_MARKERS = (
    "/image",
    "/media",
    "/video",
    "/audio",
    "/oleObject",
    "/package",
)
POWERPOINT_EXTENSIONS = {".pptx", ".pptm"}


@dataclass
class SlideArchive:
    number: int
    screenshot_link: str = ""
    text_blocks: list[str] = field(default_factory=list)
    notes: str = ""
    media_links: list[str] = field(default_factory=list)


@dataclass
class ConversionResult:
    source: Path
    output_dir: Path
    markdown_path: Path | None = None
    title: str = ""
    slide_count: int = 0
    screenshot_count: int = 0
    media_count: int = 0
    status: str = "ok"
    error: str = ""
    render_warning: str = ""


def slugify(value: str, fallback: str = "presentation") -> str:
    """Return a conservative filename-safe slug while preserving readability."""
    value = value.strip()
    value = re.sub(r"[^\w\s.-]+", "", value, flags=re.UNICODE)
    value = re.sub(r"[\s_]+", "_", value)
    value = re.sub(r"_+", "_", value).strip("._- ")
    return value or fallback


def unique_output_dir(base_output: Path, source_stem: str, overwrite: bool) -> Path:
    folder = base_output / slugify(source_stem)
    if overwrite or not folder.exists():
        return folder

    counter = 2
    while True:
        candidate = base_output / f"{slugify(source_stem)}_{counter}"
        if not candidate.exists():
            return candidate
        counter += 1


def find_pptx_files(input_path: Path, recursive: bool) -> list[Path]:
    if input_path.is_file():
        return [input_path] if input_path.suffix.lower() in POWERPOINT_EXTENSIONS else []

    pattern = "**/*" if recursive else "*"
    return sorted(
        p
        for p in input_path.glob(pattern)
        if p.is_file() and p.suffix.lower() in POWERPOINT_EXTENSIONS
    )


def shape_text_blocks(shape) -> list[str]:
    blocks: list[str] = []

    if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
        for child in shape.shapes:
            blocks.extend(shape_text_blocks(child))

    if getattr(shape, "has_text_frame", False):
        paragraphs: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                paragraphs.append(text)
        if paragraphs:
            blocks.append("\n".join(paragraphs))

    if getattr(shape, "has_table", False):
        rows: list[str] = []
        for row in shape.table.rows:
            cells = []
            for cell in row.cells:
                cell_text = " ".join(
                    paragraph.text.strip()
                    for paragraph in cell.text_frame.paragraphs
                    if paragraph.text.strip()
                )
                cells.append(cell_text)
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append("\n".join(rows))

    return blocks


def slide_text_blocks(slide) -> list[str]:
    blocks: list[str] = []
    seen: set[str] = set()
    for shape in slide.shapes:
        for block in shape_text_blocks(shape):
            normalized = normalize_text(block)
            if normalized and normalized not in seen:
                seen.add(normalized)
                blocks.append(normalized)
    return blocks


def slide_notes(slide) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        if text_frame is None:
            return ""
        return normalize_text(text_frame.text)
    except Exception:
        return ""


def normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.split("\n")]
    while lines and not lines[0]:
        lines.pop(0)
    while lines and not lines[-1]:
        lines.pop()
    return "\n".join(lines)


def presentation_title(prs, source: Path) -> str:
    try:
        if prs.slides and prs.slides[0].shapes.title:
            title = normalize_text(prs.slides[0].shapes.title.text)
            if title:
                return title.splitlines()[0]
    except Exception:
        pass
    return source.stem


def relationship_sort_key(rel: ET.Element) -> tuple[int, str]:
    rid = rel.attrib.get("Id", "")
    match = re.search(r"(\d+)$", rid)
    return (int(match.group(1)) if match else 0, rid)


def relationship_target_to_zip_path(slide_rels_path: str, target: str) -> str:
    base = PurePosixPath(slide_rels_path).parent.parent
    resolved = base.joinpath(target)
    parts: list[str] = []
    for part in resolved.parts:
        if part in ("", "."):
            continue
        if part == "..":
            if parts:
                parts.pop()
            continue
        parts.append(part)
    return "/".join(parts)


def extract_slide_media(source: Path, images_dir: Path) -> dict[int, list[str]]:
    """Extract slide-linked media from the pptx zip and return slide-number links."""
    slide_links: dict[int, list[str]] = {}
    if not zipfile.is_zipfile(source):
        return slide_links

    with zipfile.ZipFile(source) as archive:
        names = set(archive.namelist())
        rel_paths = sorted(
            name
            for name in names
            if re.fullmatch(r"ppt/slides/_rels/slide\d+\.xml\.rels", name)
        )

        for rel_path in rel_paths:
            slide_match = re.search(r"slide(\d+)\.xml\.rels$", rel_path)
            if not slide_match:
                continue
            slide_number = int(slide_match.group(1))

            try:
                root = ET.fromstring(archive.read(rel_path))
            except ET.ParseError:
                continue

            media_index = 1
            for rel in sorted(root.findall(f"{REL_NS}Relationship"), key=relationship_sort_key):
                rel_type = rel.attrib.get("Type", "")
                target = rel.attrib.get("Target", "")
                if not target or target.startswith("http"):
                    continue
                if not any(marker in rel_type for marker in MEDIA_REL_MARKERS):
                    continue

                zip_media_path = relationship_target_to_zip_path(rel_path, target)
                if zip_media_path not in names:
                    continue

                original_name = PurePosixPath(zip_media_path).name
                suffix = Path(original_name).suffix or ".bin"
                safe_name = (
                    f"slide_{slide_number:03d}_image_{media_index:03d}"
                    f"{suffix.lower()}"
                )
                destination = images_dir / safe_name
                destination.write_bytes(archive.read(zip_media_path))
                slide_links.setdefault(slide_number, []).append(f"images/{safe_name}")
                media_index += 1

    return slide_links


def first_available_command(names: Iterable[str]) -> str | None:
    for name in names:
        command = shutil.which(name)
        if command:
            return command
    return None


def render_slide_screenshots(
    source: Path,
    images_dir: Path,
    expected_slide_count: int,
    dpi: int = 160,
) -> tuple[dict[int, str], str]:
    """Render full-slide screenshots via LibreOffice plus pdftoppm."""
    office = first_available_command(("soffice", "libreoffice"))
    pdftoppm = first_available_command(("pdftoppm",))

    if not office:
        return {}, "Slide screenshots skipped: LibreOffice/soffice was not found."
    if not pdftoppm:
        return {}, "Slide screenshots skipped: pdftoppm was not found."

    screenshot_links: dict[int, str] = {}
    with tempfile.TemporaryDirectory(prefix="pptx_md_render_") as tmp_name:
        tmp_dir = Path(tmp_name)
        convert_cmd = [
            office,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(tmp_dir),
            str(source),
        ]
        converted = subprocess.run(
            convert_cmd,
            capture_output=True,
            text=True,
            timeout=300,
            check=False,
        )
        if converted.returncode != 0:
            details = (converted.stderr or converted.stdout).strip()
            return {}, f"Slide screenshots skipped: LibreOffice PDF export failed. {details}"

        pdf_path = tmp_dir / f"{source.stem}.pdf"
        if not pdf_path.exists():
            pdf_candidates = sorted(tmp_dir.glob("*.pdf"))
            if not pdf_candidates:
                return {}, "Slide screenshots skipped: LibreOffice did not create a PDF."
            pdf_path = pdf_candidates[0]

        image_prefix = tmp_dir / "slide_render"
        render_cmd = [
            pdftoppm,
            "-png",
            "-r",
            str(dpi),
            str(pdf_path),
            str(image_prefix),
        ]
        rendered = subprocess.run(
            render_cmd,
            capture_output=True,
            text=True,
            timeout=300,
            check=False,
        )
        if rendered.returncode != 0:
            details = (rendered.stderr or rendered.stdout).strip()
            return {}, f"Slide screenshots skipped: pdftoppm render failed. {details}"

        rendered_pages = sorted(
            tmp_dir.glob("slide_render-*.png"),
            key=lambda path: int(re.search(r"-(\d+)\.png$", path.name).group(1)),
        )

        for slide_number, rendered_page in enumerate(rendered_pages, start=1):
            if expected_slide_count and slide_number > expected_slide_count:
                break
            safe_name = f"slide_{slide_number:03d}_screenshot.png"
            destination = images_dir / safe_name
            shutil.copy2(rendered_page, destination)
            screenshot_links[slide_number] = f"images/{safe_name}"

    if expected_slide_count and len(screenshot_links) < expected_slide_count:
        warning = (
            "Slide screenshots partially rendered: "
            f"{len(screenshot_links)} of {expected_slide_count} slides."
        )
    else:
        warning = ""
    return screenshot_links, warning


def markdown_link(path_text: str) -> str:
    label = Path(path_text).name
    href = path_text.replace("\\", "/").replace(" ", "%20")
    return f"[{label}]({href})"


def write_markdown(
    result: ConversionResult,
    slides: list[SlideArchive],
    include_notes: bool,
) -> Path:
    md_path = result.output_dir / f"{slugify(result.source.stem)}.md"
    converted = datetime.now().astimezone().isoformat(timespec="seconds")

    lines: list[str] = [
        f"# {result.title}",
        "",
        f"Source file: {result.source}",
        f"Converted: {converted}",
        "",
    ]

    for slide in slides:
        lines.extend([f"## Slide {slide.number}", ""])

        lines.extend(["### Slide screenshot", ""])
        if slide.screenshot_link:
            lines.extend([f"- {markdown_link(slide.screenshot_link)}", ""])
        else:
            lines.extend(["_No slide screenshot rendered._", ""])

        lines.extend(["### Text", ""])
        if slide.text_blocks:
            for block in slide.text_blocks:
                lines.extend([block, ""])
        else:
            lines.extend(["_No extractable slide text found._", ""])

        if include_notes:
            lines.extend(["### Speaker notes", ""])
            lines.extend([slide.notes or "_No speaker notes found._", ""])

        lines.extend(["### Images", ""])
        if slide.media_links:
            for link in slide.media_links:
                lines.append(f"- {markdown_link(link)}")
            lines.append("")
        else:
            lines.extend(["_No slide-linked images or media extracted._", ""])

    md_path.write_text("\n".join(lines), encoding="utf-8", newline="\n")
    return md_path


def convert_presentation(
    source: Path,
    output_base: Path,
    overwrite: bool,
    extract_images: bool,
    slide_screenshots: bool,
    include_notes: bool,
) -> ConversionResult:
    result = ConversionResult(source=source, output_dir=output_base)

    if not source.exists() or source.stat().st_size == 0:
        result.status = "failed"
        result.error = "File is missing or empty."
        return result

    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise SystemExit(
            "python-pptx is required. Install it with: python -m pip install python-pptx"
        ) from exc

    result.output_dir = unique_output_dir(output_base, source.stem, overwrite)
    if result.output_dir.exists() and overwrite:
        for existing in result.output_dir.glob("*"):
            if existing.is_file():
                existing.unlink()
            elif existing.name == "images":
                for image_file in existing.glob("*"):
                    if image_file.is_file():
                        image_file.unlink()
    result.output_dir.mkdir(parents=True, exist_ok=True)
    images_dir = result.output_dir / "images"
    images_dir.mkdir(exist_ok=True)

    try:
        prs = Presentation(str(source))
        result.title = presentation_title(prs, source)
        result.slide_count = len(prs.slides)

        screenshot_by_slide: dict[int, str] = {}
        if slide_screenshots:
            screenshot_by_slide, result.render_warning = render_slide_screenshots(
                source=source,
                images_dir=images_dir,
                expected_slide_count=result.slide_count,
            )

        media_by_slide = extract_slide_media(source, images_dir) if extract_images else {}

        slides: list[SlideArchive] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_archive = SlideArchive(number=idx)
            slide_archive.screenshot_link = screenshot_by_slide.get(idx, "")
            slide_archive.text_blocks = slide_text_blocks(slide)
            slide_archive.notes = slide_notes(slide)
            slide_archive.media_links = media_by_slide.get(idx, [])
            slides.append(slide_archive)

        result.screenshot_count = sum(1 for slide in slides if slide.screenshot_link)
        result.media_count = sum(len(slide.media_links) for slide in slides)
        result.markdown_path = write_markdown(result, slides, include_notes=include_notes)

        metadata = {
            "source": str(source),
            "markdown": str(result.markdown_path),
            "title": result.title,
            "slide_count": result.slide_count,
            "screenshot_count": result.screenshot_count,
            "media_count": result.media_count,
            "render_warning": result.render_warning,
            "converted": datetime.now().astimezone().isoformat(timespec="seconds"),
        }
        (result.output_dir / "conversion_metadata.json").write_text(
            json.dumps(metadata, indent=2, ensure_ascii=False),
            encoding="utf-8",
            newline="\n",
        )
    except Exception as exc:
        result.status = "failed"
        result.error = str(exc)
        error_md = result.output_dir / f"{slugify(source.stem)}.md"
        error_md.write_text(
            f"# {source.stem}\n\nSource file: {source}\n\nConversion failed: {exc}\n",
            encoding="utf-8",
            newline="\n",
        )
        result.markdown_path = error_md

    return result


def write_index(output_base: Path, results: Iterable[ConversionResult]) -> Path:
    index_path = output_base / "index.md"
    converted = datetime.now().astimezone().isoformat(timespec="seconds")
    lines = [
        "# Converted PowerPoint Archives",
        "",
        f"Generated: {converted}",
        "",
        "| Presentation | Slides | Screenshots | Media | Status |",
        "| --- | ---: | ---: | ---: | --- |",
    ]

    for result in results:
        if result.markdown_path:
            rel = result.markdown_path.relative_to(output_base).as_posix()
            title = result.title or result.source.stem
            presentation = f"[{title}]({rel.replace(' ', '%20')})"
        else:
            presentation = result.source.name
        status = result.status if result.status == "ok" else f"failed: {result.error}"
        if result.status == "ok" and result.render_warning:
            status = f"ok; {result.render_warning}"
        lines.append(
            f"| {presentation} | {result.slide_count} | {result.screenshot_count} | "
            f"{result.media_count} | {status} |"
        )

    index_path.write_text("\n".join(lines) + "\n", encoding="utf-8", newline="\n")
    return index_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Convert PowerPoint .pptx/.pptm files into Markdown archives with extracted "
            "slide-linked images/media."
        )
    )
    parser.add_argument("--input", required=True, help="Path to a .pptx/.pptm file or folder.")
    parser.add_argument("--output", required=True, help="Path to the output folder.")
    parser.add_argument("--recursive", action="store_true", help="Search folders recursively.")
    parser.add_argument("--overwrite", action="store_true", help="Overwrite matching output folders.")
    parser.add_argument(
        "--extract-images",
        action="store_true",
        default=True,
        help="Extract slide-linked images/media. Enabled by default.",
    )
    parser.add_argument(
        "--slide-screenshots",
        dest="slide_screenshots",
        action="store_true",
        default=True,
        help="Render and link one full-slide screenshot per slide. Enabled by default.",
    )
    parser.add_argument(
        "--no-slide-screenshots",
        dest="slide_screenshots",
        action="store_false",
        help="Skip full-slide screenshot rendering.",
    )
    parser.add_argument("--include-notes", action="store_true", help="Include speaker notes.")
    parser.add_argument(
        "--make-index",
        action="store_true",
        help="Write index.md listing converted presentations. Index output is always enabled.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    input_path = Path(args.input).expanduser().resolve()
    output_base = Path(args.output).expanduser().resolve()
    output_base.mkdir(parents=True, exist_ok=True)

    pptx_files = find_pptx_files(input_path, recursive=args.recursive)
    if not pptx_files:
        print(f"No .pptx/.pptm files found in: {input_path}", file=sys.stderr)
        return 1

    results: list[ConversionResult] = []
    for source in pptx_files:
        result = convert_presentation(
            source=source,
            output_base=output_base,
            overwrite=args.overwrite,
            extract_images=args.extract_images,
            slide_screenshots=args.slide_screenshots,
            include_notes=args.include_notes,
        )
        results.append(result)
        if result.status == "ok":
            print(f"Converted: {source} -> {result.markdown_path}")
        else:
            print(f"Failed: {source} ({result.error})", file=sys.stderr)

    index_path = write_index(output_base, results)
    print(f"Index: {index_path}")

    failed = sum(1 for result in results if result.status != "ok")
    return 1 if failed and failed == len(results) else 0


if __name__ == "__main__":
    raise SystemExit(main())
